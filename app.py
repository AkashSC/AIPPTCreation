import io
import os
import re
import json
import streamlit as st
import pdfplumber
from docx import Document
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from groq import Groq

# ------------------------------
# Config
# ------------------------------
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
DEFAULT_MODEL = "llama3-8b-8192"
client = Groq(api_key=GROQ_API_KEY)

# ------------------------------
# Utility: text extraction
# ------------------------------
def extract_text_from_pdf(file_bytes: bytes) -> str:
    text_parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            txt = page.extract_text()
            if txt:
                text_parts.append(txt)
    return "\n\n".join(text_parts).strip()

def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    paras = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n\n".join(paras).strip()

def extract_text_from_txt(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="replace").strip()

def extract_text(uploaded_file):
    data = uploaded_file.read()
    name = uploaded_file.name.lower()
    if name.endswith(".pdf"):
        return extract_text_from_pdf(data)
    elif name.endswith(".docx") or name.endswith(".doc"):
        return extract_text_from_docx(data)
    else:
        return extract_text_from_txt(data)

# ------------------------------
# Local summarizer fallback
# ------------------------------
def simple_local_summary(text: str, max_sentences: int = 4) -> str:
    t = re.sub(r"\s+", " ", text).strip()
    sentences = re.split(r'(?<=[.!?])\s+', t)
    return " ".join(sentences[:max_sentences]) or (t[:200] + ("..." if len(t) > 200 else ""))

# ------------------------------
# Parse design prompt into basic style (fallback)
# ------------------------------
COLOR_MAP = {
    "blue": "#003366", "dark blue": "#003366", "black": "#000000", "white": "#FFFFFF",
    "green": "#008000", "red": "#FF0000", "yellow": "#FFCC00", "gray": "#808080",
    "dark": "#333333", "light": "#F8F8F8", "orange": "#FF8C00", "purple": "#800080"
}
FONT_OPTIONS = ["Arial", "Calibri", "Times New Roman", "Helvetica", "Comic Sans MS", "Verdana"]

def parse_design_prompt(prompt: str):
    prompt_l = (prompt or "").lower()
    style = {"background_color": "#FFFFFF", "font": "Arial", "font_size": 14, "font_color": "#000000", "emoji_in_bullets": False}

    # hex color
    found_hex = re.search(r'#([0-9a-fA-F]{6})', prompt)
    if found_hex:
        style["background_color"] = f"#{found_hex.group(1)}"

    # color words
    for name, hx in COLOR_MAP.items():
        if name in prompt_l:
            style["background_color"] = hx
            # choose readable default font color
            style["font_color"] = "#FFFFFF" if name in ("blue","dark blue","black","dark","purple") else "#000000"
            break

    # font
    for f in FONT_OPTIONS:
        if f.lower() in prompt_l:
            style["font"] = f
            break

    # font size hints
    if "large" in prompt_l or "big" in prompt_l or "title large" in prompt_l:
        style["font_size"] = 20
    if "small" in prompt_l or "compact" in prompt_l:
        style["font_size"] = 12
    if m := re.search(r'font ?size ?[:= ]?(\d{2})', prompt_l):
        try:
            style["font_size"] = int(m.group(1))
        except:
            pass

    # emojis in bullets
    if "emoji" in prompt_l or "emojis" in prompt_l or "smiley" in prompt_l:
        style["emoji_in_bullets"] = True

    return style

# ------------------------------
# Helpers: JSON extraction & slide parsing
# ------------------------------
def extract_style_json_from_text(s: str):
    # 1) try tags <STYLE_JSON>...</STYLE_JSON>
    m = re.search(r'<STYLE_JSON>(.*?)</STYLE_JSON>', s, re.DOTALL | re.IGNORECASE)
    if m:
        try:
            return json.loads(m.group(1).strip())
        except:
            pass
    # 2) try to find last balanced JSON object
    for start in range(len(s)-1, -1, -1):
        if s[start] == '{':
            depth = 0
            for end in range(start, len(s)):
                if s[end] == '{':
                    depth += 1
                elif s[end] == '}':
                    depth -= 1
                if depth == 0:
                    candidate = s[start:end+1]
                    try:
                        return json.loads(candidate)
                    except:
                        break
    return None

def parse_slides_from_output(output: str):
    slides = []
    # try structured pattern: "Slide Title: ..." or "Title: ..."
    pattern = re.compile(r'(?:Slide Title:|Title:)\s*(.+?)(?:\n|$)([\s\S]*?)(?=(?:Slide Title:|Title:)|$)', re.IGNORECASE)
    matches = pattern.findall(output)
    if matches:
        for title, body in matches:
            title = title.strip()
            # bullets are lines starting with -, •, * or numbered
            bullets = []
            for line in body.splitlines():
                line = line.strip()
                if not line:
                    continue
                if re.match(r'^[-•\*\d\)]\s+', line):
                    bullets.append(re.sub(r'^[-•\*\d\)\.]+\s*', '', line).strip())
            if not bullets:
                # fallback: create bullets from body using local summary
                summary = simple_local_summary(body, max_sentences=4)
                bullets = re.split(r'(?<=[.!?])\s+', summary)
            slides.append({"title": title, "bullets": bullets})
        return slides

    # if no structured titles found, try markdown-style "# Title"
    md_matches = re.findall(r'^\s*#\s*(.+)$', output, re.MULTILINE)
    if md_matches:
        # split roughly
        parts = re.split(r'^\s*#\s*.+$', output, flags=re.MULTILINE)
        for i, t in enumerate(md_matches):
            body = parts[i+1] if i+1 < len(parts) else ""
            bullets = [l.strip().lstrip('-•* ').strip() for l in body.splitlines() if l.strip().startswith('-')]
            if not bullets:
                summary = simple_local_summary(body or t, max_sentences=4)
                bullets = re.split(r'(?<=[.!?])\s+', summary)
            slides.append({"title": t.strip(), "bullets": bullets})
        return slides

    # final fallback: chunk the text and create slides via local summarizer
    cleaned = re.sub(r'\s+', ' ', output).strip()
    if not cleaned:
        return []
    chunk_size = 1200
    chunks = [cleaned[i:i+chunk_size] for i in range(0, len(cleaned), chunk_size)]
    for i, ch in enumerate(chunks, start=1):
        summary = simple_local_summary(ch, max_sentences=4)
        bullets = re.split(r'(?<=[.!?])\s+', summary)
        title = bullets[0] if bullets else f"Part {i}"
        slides.append({"title": title[:60], "bullets": bullets})
    return slides

# ------------------------------
# LLM call: ask for slides + style JSON (robust)
# ------------------------------
def summarize_and_style_with_groq(text: str, design_prompt: str, model: str = DEFAULT_MODEL, max_chunk_chars: int = 3000):
    """
    Returns (slides, style_json or None, raw_output, used_groq_bool, error_message_or_none)
    """
    # keep prompts small: chunk text if huge
    text_for_prompt = text if len(text) <= max_chunk_chars else text[:max_chunk_chars] + "\n\n[TRUNCATED]"
    prompt = f"""
You are a helpful presentation designer. Summarize the supplied document into presentation slides and produce a JSON with style settings.

Requirements:
1) For each slide, produce a title and 4-5 concise bullet points (one per line).
2) Follow these design instructions: {design_prompt}
3) After the slides, include a JSON block enclosed by <STYLE_JSON>...</STYLE_JSON> exactly (only the JSON inside).
   JSON keys should include (if possible): background_color (hex like #RRGGBB), font (string), font_size (number), font_color (hex).
4) Keep the slides short and presentation-friendly.

Example output:
Slide Title: Example Slide
- Bullet one
- Bullet two

STYLE_JSON:
<STYLE_JSON>{"background_color":"#003366","font":"Calibri","font_size":18,"font_color":"#FFFFFF"}</STYLE_JSON>

Document:
{text_for_prompt}
"""
    try:
        chat = client.chat.completions.create(
            model=model,
            messages=[{"role": "system", "content": "You are a presentation designer."},
                      {"role": "user", "content": prompt}],
            temperature=0.3,
            max_tokens=800
        )
        raw = chat.choices[0].message.content
        slides = parse_slides_from_output(raw)
        style_json = extract_style_json_from_text(raw)
        return slides, style_json, raw, True, None
    except Exception as e:
        return [], None, "", False, str(e)

# ------------------------------
# PPT construction
# ------------------------------
def hex_to_rgb_obj(hex_color: str):
    if not hex_color:
        hex_color = "#FFFFFF"
    hex_color = hex_color.strip()
    if not hex_color.startswith("#"):
        # try color map
        hex_color = COLOR_MAP.get(hex_color.lower(), "#FFFFFF")
    try:
        h = hex_color.lstrip("#")
        if len(h) != 6:
            raise ValueError("Invalid length")
        r = int(h[0:2], 16); g = int(h[2:4], 16); b = int(h[4:6], 16)
        return RGBColor(r, g, b)
    except:
        return RGBColor(255, 255, 255)

def make_ppt(slides, style):
    # normalize style keys
    style = style or {}
    bg = style.get("background_color", "#FFFFFF")
    font = style.get("font", "Arial")
    try:
        font_size = int(style.get("font_size", 14))
    except:
        font_size = 14
    font_color = style.get("font_color", "#000000")
    emoji_flag = style.get("emoji_in_bullets", False)

    prs = Presentation()
    # Title slide
    ts = prs.slides.add_slide(prs.slide_layouts[0])
    ts.shapes.title.text = "Auto-generated Presentation"
    if ts.placeholders:
        try:
            ts.placeholders[1].text = "Generated by Groq + Agentic pipeline"
        except:
            pass

    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        # background
        try:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb_obj(bg)
        except Exception:
            pass

        # title
        try:
            title_shape = slide.shapes.title
            title_shape.text = s.get("title", "")[:250]
            title_shape.text_frame.paragraphs[0].font.name = font
            title_shape.text_frame.paragraphs[0].font.size = Pt(max(font_size + 4, 14))
            title_shape.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb_obj(font_color)
        except Exception:
            pass

        # bullets
        try:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for b in s.get("bullets", []):
                text_b = b
                if emoji_flag:
                    text_b = "• " + text_b  # user prompt wanted emojis, LLM might already add them
                p = tf.add_paragraph()
                p.text = text_b
                p.level = 0
                p.font.name = font
                p.font.size = Pt(font_size)
                p.font.color.rgb = hex_to_rgb_obj(font_color)
        except Exception:
            # best-effort: create a single textbox if placeholders fail
            try:
                left = top = width = height = Inches(1)
                tx = slide.shapes.add_textbox(left, top, width, height)
                tf = tx.text_frame
                tf.text = "\n".join(s.get("bullets", []))[:1000]
            except:
                pass

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()

# ------------------------------
# Streamlit UI
# ------------------------------
st.set_page_config(page_title="Agentic PPT generator", layout="wide")
st.title("📄 ➜ 🖥️ Multi-doc → PPT (Agentic, design-prompts applied)")

st.markdown("Upload documents and provide free-form design instructions. The model will summarize and we will apply the design to the PPT.")

# Design prompt input + examples (below box)
design_prompt = st.text_area("Design instructions (free-form). Examples below the box:", value="Blue gradient background, bold titles, Calibri, large font, add emojis to bullets")
st.markdown("""
**Design prompt examples (copy/paste or edit):**
- `Blue gradient background, Calibri, white bold titles, font size 20, add emojis to bullets`
- `Minimalist white background, Helvetica, small font, black text`
- `Dark theme (#0f172a), Sans-serif, font size 18, yellow bullet points`
- `Playful: use Comic Sans MS, add emojis, large font, pastel background #FFDDE1`
""")

st.write("---")
files = st.file_uploader("Upload PDF / DOCX / TXT (multiple)", accept_multiple_files=True, type=["pdf","docx","txt"])
model_choice = st.selectbox("Groq model", ["llama3-8b-8192","gemma2-9b-it","mixtral-8x7b"])

if files and st.button("Generate PPT"):
    all_slides = []
    # base style from design prompt
    global_style = parse_design_prompt(design_prompt)

    any_groq_used = False
    groq_errors = []

    for f in files:
        st.info(f"Processing {f.name} ...")
        text = extract_text(f)
        if not text:
            st.warning(f"No extractable text in {f.name}; skipping.")
            continue

        slides, style_json, raw, used, err = summarize_and_style_with_groq(text, design_prompt, model=model_choice)
        if used:
            any_groq_used = True
            st.success(f"Groq summarization succeeded for {f.name}. Parsed {len(slides)} slides.")
            if style_json:
                # merge style_json into global_style (json overrides prompt-based parse)
                try:
                    # normalize keys lower-case
                    for k, v in style_json.items():
                        global_style[k] = v
                except Exception:
                    pass
        else:
            groq_errors.append((f.name, err))
            st.warning(f"Groq failed for {f.name}: {err}. Falling back to local summarizer.")
            # fallback: create 1 or more slides via local summarizer
            # chunk text and produce slides
            chunk_size = 1800
            chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
            for i, ch in enumerate(chunks, start=1):
                summ = simple_local_summary(ch, max_sentences=4)
                bullets = re.split(r'(?<=[.!?])\s+', summ)
                title = bullets[0] if bullets else f"{f.name} - Part {i}"
                all_slides.append({"title": title[:60], "bullets": bullets})
            continue

        # if we got slides from Groq, append them
        if slides:
            all_slides.extend(slides)
        else:
            # as last fallback, use local summary for the file
            summ = simple_local_summary(text, max_sentences=4)
            bullets = re.split(r'(?<=[.!?])\s+', summ)
            title = bullets[0] if bullets else f"{f.name}"
            all_slides.append({"title": title[:60], "bullets": bullets})

    # if Groq never returned style JSON, try to parse the user's prompt
    # ensure some keys present
    # also allow parse_design_prompt to detect emoji instruction
    parsed_from_prompt = parse_design_prompt(design_prompt)
    # merge: existing global_style already contains any LLM JSON overrides; fill missing from parsed prompt
    for k, v in parsed_from_prompt.items():
        if k not in global_style or global_style.get(k) is None:
            global_style[k] = v

    if not all_slides:
        st.error("No slides could be generated from the uploaded files.")
    else:
        pptx_bytes = make_ppt(all_slides, global_style)
        st.success(f"Generated PPT with {len(all_slides)} slides.")
        st.download_button("⬇️ Download PPTX", pptx_bytes, file_name="auto_presentation.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    if groq_errors:
        st.error("Groq errors occurred for some files (fallback applied):")
        for name, err in groq_errors:
            st.write(f"- {name}: {err}")
